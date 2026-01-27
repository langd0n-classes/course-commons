#!/usr/bin/env python3
"""
Extract text from PPTX files into .txt sidecar files.
"""

from __future__ import annotations

import argparse
import glob
import json
import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError as exc:  # pragma: no cover
    raise SystemExit(
        "Missing dependency: python-pptx. Install with: pip install -r requirements.txt"
    ) from exc


def extract_text(path: Path, include_notes: bool, include_layout: bool) -> list[str]:
    prs = Presentation(str(path))
    out: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"=== Slide {i} ===")
        if include_layout:
            layout_name = getattr(slide.slide_layout, "name", "")
            title_shape = getattr(slide.shapes, "title", None)
            title_text = title_shape.text.strip() if title_shape and title_shape.text else ""
            if layout_name:
                out.append(f"[Layout] {layout_name}")
            if title_text:
                out.append(f"[Title] {title_text}")
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                out.append(text.strip())
        if include_notes:
            notes = ""
            try:
                notes = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                notes = ""
            if notes:
                out.append("[Notes]")
                out.append(notes)
        out.append("")
    return out


def write_output(src: Path, lines: list[str], out_dir: Path | None) -> Path:
    if out_dir:
        out_dir.mkdir(parents=True, exist_ok=True)
        dst = out_dir / (src.stem + ".txt")
    else:
        dst = src.with_suffix(".txt")
    dst.write_text("\n".join(lines), encoding="utf-8")
    return dst


def write_json_output(src: Path, slides: list[dict], out_dir: Path | None) -> Path:
    if out_dir:
        out_dir.mkdir(parents=True, exist_ok=True)
        dst = out_dir / (src.stem + ".json")
    else:
        dst = src.with_suffix(".json")
    dst.write_text(json.dumps(slides, indent=2), encoding="utf-8")
    return dst


def parse_args(argv: list[str]) -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description=(
            "Extract text from PPTX files and write .txt outputs for GenAI-friendly parsing."
        )
    )
    parser.add_argument(
        "paths",
        nargs="*",
        default=["*.pptx"],
        help="Input PPTX files or glob patterns (default: *.pptx).",
    )
    parser.add_argument(
        "--out-dir",
        help="Output directory for .txt files (default: next to PPTX).",
    )
    parser.add_argument(
        "--overwrite",
        action="store_true",
        help="Overwrite existing .txt files (default: skip).",
    )
    parser.add_argument(
        "--no-notes",
        action="store_true",
        help="Do not include speaker notes (notes included by default).",
    )
    parser.add_argument(
        "--include-layout",
        action="store_true",
        help="Include slide layout name and title when present.",
    )
    parser.add_argument(
        "--json",
        action="store_true",
        help="Write JSON output instead of plain text.",
    )
    return parser.parse_args(argv)


def main(argv: list[str]) -> int:
    args = parse_args(argv)
    patterns = args.paths
    include_notes = not args.no_notes

    matched: list[Path] = []
    for pattern in patterns:
        expanded = [Path(p) for p in glob.glob(pattern)]
        matched.extend(expanded)

    if not matched:
        print("No PPTX files found. Provide files or glob patterns.")
        return 2

    out_dir = Path(args.out_dir) if args.out_dir else None
    failures = 0

    for path in matched:
        if not path.exists():
            print(f"Skip (not found): {path}")
            failures += 1
            continue
        if path.suffix.lower() != ".pptx":
            print(f"Skip (not .pptx): {path}")
            continue

        try:
            if args.json:
                slides = []
                prs = Presentation(str(path))
                for i, slide in enumerate(prs.slides, 1):
                    entry = {"slide": i, "texts": []}
                    if args.include_layout:
                        entry["layout"] = getattr(slide.slide_layout, "name", "")
                        title_shape = getattr(slide.shapes, "title", None)
                        entry["title"] = title_shape.text.strip() if title_shape and title_shape.text else ""
                    for shape in slide.shapes:
                        text = getattr(shape, "text", "")
                        if text and text.strip():
                            entry["texts"].append(text.strip())
                    if include_notes:
                        try:
                            notes = slide.notes_slide.notes_text_frame.text.strip()
                        except Exception:
                            notes = ""
                        if notes:
                            entry["notes"] = notes
                    slides.append(entry)
                dst = (out_dir / (path.stem + ".json")) if out_dir else path.with_suffix(".json")
                if dst.exists() and not args.overwrite:
                    print(f"Skip (exists): {dst}")
                    continue
                dst = write_json_output(path, slides, out_dir)
            else:
                lines = extract_text(path, include_notes, args.include_layout)
                dst = (out_dir / (path.stem + ".txt")) if out_dir else path.with_suffix(".txt")
                if dst.exists() and not args.overwrite:
                    print(f"Skip (exists): {dst}")
                    continue
                dst = write_output(path, lines, out_dir)
            print(f"Wrote: {dst}")
        except Exception as exc:  # pragma: no cover
            print(f"Error processing {path}: {exc}")
            failures += 1

    return 1 if failures else 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
